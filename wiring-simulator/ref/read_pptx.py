# -*- coding: utf-8 -*-
"""記入済みの PowerPoint から、部品の位置と配線のルートを数値で読み取る。

  使い方: python wiring-simulator/ref/read_pptx.py <返ってきた.pptx>

  ⚠️縮尺は「画像図形の実際の位置と大きさ」から毎回逆算する（対応表は画像の名前に
    CARVIEW|view|m0|m1|v0|v1 の形で埋めてある）＝図を動かされても拡大されても正しく読める。
  ⚠️部品の位置は【図形の中央】。パレット（図の外）に置きっぱなしの物は「未配置」として弾く。
  ⚠️配線はフリーフォームの折れ点を拾い、両端にいちばん近い部品名を当てる。
"""
import sys, json
from pptx import Presentation
from pptx.util import Emu

SRC = sys.argv[1] if len(sys.argv) > 1 else '_handwrite_FIAT500F.pptx'
NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'


def mm(v):
    return Emu(int(v)).mm


def mapper(pic):
    """画像図形から「スライド上のmm → 車の座標(m)」の換算関数を作る。"""
    _, view, m0, m1, v0, v1 = pic.name.split('|')
    m0, m1, v0, v1 = map(float, (m0, m1, v0, v1))
    L, T = mm(pic.left), mm(pic.top)
    W, H = mm(pic.width), mm(pic.height)

    def to_car(x, y):
        return (m0 + (x - L) / W * (m1 - m0), v0 + (y - T) / H * (v1 - v0))

    def inside(x, y, slack=3.0):
        return L - slack <= x <= L + W + slack and T - slack <= y <= T + H + slack
    return view, to_car, inside


def freeform_points(sp):
    """フリーフォームの折れ点をスライド上のmmで返す。custGeom を直接読む。"""
    el = sp._element
    g = el.find('.//' + NS + 'custGeom')
    if g is None:
        return []
    # ⚠️'.//a:ext' だけで探すと extLst の中の別物（uri付き）を拾う＝必ず xfrm 配下から取る
    xf = el.find('.//' + NS + 'xfrm')
    if xf is None:
        return []
    ext, off = xf.find(NS + 'ext'), xf.find(NS + 'off')
    if ext is None or off is None:
        return []
    cx, cy = float(ext.get('cx')), float(ext.get('cy'))
    ox, oy = float(off.get('x')), float(off.get('y'))
    pts = []
    for path in g.iter(NS + 'path'):
        w, h = float(path.get('w') or 0), float(path.get('h') or 0)
        if w <= 0 or h <= 0:
            continue
        for node in path:
            for p in node.iter(NS + 'pt'):
                x = ox + float(p.get('x')) / w * cx
                y = oy + float(p.get('y')) / h * cy
                pts.append((mm(x), mm(y)))
    return pts


def connector_points(sp):
    """コネクタ（直線・カギ線）の通り道をスライド上のmmで返す。
       ⚠️フリーフォームだけ見ていると、直線コネクタで引かれた線を丸ごと取りこぼす。"""
    el = sp._element
    pg = el.find('.//' + NS + 'prstGeom')
    if pg is None or pg.get('prst') not in (
            'line', 'straightConnector1', 'bentConnector2', 'bentConnector3'):
        return []
    xf = el.find('.//' + NS + 'xfrm')
    if xf is None:
        return []
    off, ext = xf.find(NS + 'off'), xf.find(NS + 'ext')
    if off is None or ext is None:
        return []
    x0, y0 = float(off.get('x')), float(off.get('y'))
    w, h = float(ext.get('cx')), float(ext.get('cy'))
    fh, fv = xf.get('flipH') == '1', xf.get('flipV') == '1'      # 向きは flip で表される
    ax0, ay0 = (x0 + w if fh else x0), (y0 + h if fv else y0)
    ax1, ay1 = (x0 if fh else x0 + w), (y0 if fv else y0 + h)
    if pg.get('prst') in ('line', 'straightConnector1'):
        pts = [(ax0, ay0), (ax1, ay1)]
    else:
        adj, av = .5, pg.find(NS + 'avLst')
        if av is not None:
            g = av.find(NS + 'gd')
            if g is not None and g.get('fmla'):
                adj = float(g.get('fmla').split()[-1]) / 100000.0
        mx = ax0 + (ax1 - ax0) * adj
        pts = [(ax0, ay0), (mx, ay0), (mx, ay1), (ax1, ay1)]
    return [(mm(x), mm(y)) for x, y in pts]


def rdp(pts, eps):
    """折れ点を間引く（Ramer-Douglas-Peucker）。フリーハンドは点が数百になるため。"""
    if len(pts) < 3:
        return pts
    (x0, y0), (x1, y1) = pts[0], pts[-1]
    dx, dy = x1 - x0, y1 - y0
    n = (dx * dx + dy * dy) ** .5
    best, bi = -1.0, 0
    for i in range(1, len(pts) - 1):
        x, y = pts[i]
        d = (abs(dy * x - dx * y + x1 * y0 - y1 * x0) / n) if n else             (((x - x0) ** 2 + (y - y0) ** 2) ** .5)
        if d > best:
            best, bi = d, i
    if best <= eps:
        return [pts[0], pts[-1]]
    return rdp(pts[:bi + 1], eps)[:-1] + rdp(pts[bi:], eps)


def nearest(parts, m, v):
    if not parts:
        return None, None
    d = [(((m - p['m']) ** 2 + (v - p['v']) ** 2) ** .5, p) for p in parts]
    d.sort(key=lambda t: t[0])
    return d[0][1], d[0][0]


prs = Presentation(SRC)
out = {'parts': {}, 'wires': []}
for n, sl in enumerate(prs.slides, 1):
    pic = next((s for s in sl.shapes if s.name.startswith('CARVIEW|')), None)
    if pic is None:
        continue
    view, to_car, inside = mapper(pic)
    placed, stray = [], []
    for s in sl.shapes:
        if not s.name.startswith('PART|'):
            continue
        cx, cy = mm(s.left) + mm(s.width) / 2, mm(s.top) + mm(s.height) / 2
        pid = s.name.split('|', 1)[1]
        lab = s.text_frame.text.strip() if s.has_text_frame else pid
        if not inside(cx, cy):
            stray.append(lab)
            continue
        m, v = to_car(cx, cy)
        rot = round(float(s.rotation or 0) % 360, 1)   # 向きを変えた部品（例: バッテリーの端子側）
        placed.append({'id': pid, 'label': lab, 'm': m, 'v': v, 'rot': rot})
        rec = out['parts'].setdefault(pid, {'label': lab})
        rec['label'] = lab
        if view == 'top':
            rec['m'], rec['lat'] = round(m, 3), round(v, 3)
            if rot:
                rec['rot_top'] = rot
        else:
            # 前後は上面図を正とし、側面図の値は食い違い検出のために別名で残す
            rec['h'], rec['m_side'] = round(v, 3), round(m, 3)
            if rot:
                rec['rot_side'] = rot

    print('── %d枚目（%s）' % (n, '上から' if view == 'top' else '横から'))
    if placed:
        key = '中心から(m)' if view == 'top' else '地面から(m)'
        print('   %-14s %10s %12s %8s' % ('部品', '前端から(m)', key, '回転'))
        for p in sorted(placed, key=lambda p: p['m']):
            print('   %-14s %10.3f %12.3f %8s'
                  % (p['label'], p['m'], p['v'], ('%.0f°' % p['rot']) if p['rot'] else '-'))
    if stray:
        print('   未配置（図の外に残っている）: ' + ' / '.join(stray))

    for s in sl.shapes:
        pts = freeform_points(s) or connector_points(s)
        if len(pts) < 2:
            continue
        car = rdp([to_car(x, y) for x, y in pts], 0.02)   # 20mm＝「置き場所の目安」の粒度
        a, da = nearest(placed, *car[0])
        b, db = nearest(placed, *car[-1])
        w = {'view': view,
             'from': a['id'] if a else None, 'to': b['id'] if b else None,
             'from_gap_m': round(da, 3) if da is not None else None,
             'to_gap_m': round(db, 3) if db is not None else None,
             'points': [[round(m, 3), round(v, 3)] for m, v in car]}
        out['wires'].append(w)
        print('   配線: %s → %s（元%d点→%d点・端の離れ %.02f/%.02fm）'
              % (a['label'] if a else '?', b['label'] if b else '?',
                 len(pts), len(car), da or 0, db or 0))

print()
print(json.dumps(out, ensure_ascii=False, indent=1))
