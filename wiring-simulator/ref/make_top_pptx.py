# -*- coding: utf-8 -*-
"""部品位置・配線ルートを PowerPoint 上で指定してもらうためのファイルを作る。

  紙に手書きする代わりに、パワポの図形をドラッグして置いてもらう。
  ⚠️こうする理由＝図形の座標(EMU)を読めば位置が数値で取れる。読み取りは read_pptx.py。
  ⚠️図の縮尺は「画像図形の実際の位置と大きさ」から逆算するので、
    ユーザーが図を動かしても拡大しても正しく読める（画像の名前に対応表を埋めてある）。

  出力: _手書き用_FIAT500F.pptx（repo直下・git管理外の作業ファイル）
"""
import io, re, os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DST = '_handwrite_FIAT500F.pptx'
IMGD = 'wiring-simulator/ref/_pptx_img'
DPI, SS = 300, 2                 # 印刷解像度／描くときの倍率（縮小してアンチエイリアス）
PAPER_W = 281.0                  # 図をスライド上で何mm幅で置くか

VIEWS = {
    'top':  {'src': 'wiring-img/car-top-v9.svg', 'x0': 25.0, 'base': 683.15},
    'side': {'src': 'wiring-simulator/ref/svg/fiat500f_side.svg', 'x0': 25.0, 'base': 1332.0},
}
PATHS = {k: re.findall(r'<path class="(\w+)" d="([^"]+)"',
                       io.open(v['src'], encoding='utf-8').read())
         for k, v in VIEWS.items()}

# 紙の上での線の太さ(mm)。車体のmmではない＝縮尺を変えても見え方を保つ
WIDTH_MM = {'outline': .45, 'panel': .34, 'trim': .24, 'detail': .17}
COL = {'outline': (110, 110, 110), 'panel': (125, 125, 125),
       'trim': (150, 150, 150), 'detail': (183, 183, 183)}
ZONES = [(0, 1.02, '前トランク'), (1.02, 2.05, '室内'), (2.05, 2.97, 'エンジンルーム')]
FONTP = 'C:/Windows/Fonts/meiryo.ttc'


def subpaths(d):
    out = []
    for s in d.replace('Z', ' ').split('M'):
        s = s.strip()
        if s:
            out.append([[float(a) for a in q.split(',')] for q in s.split()])
    return out


def render(view, m0, m1, v0, v1, dst):
    """view の図を PNG に描く。v0 が画像の上端、v1 が下端（top は左右m・side は高さm）。"""
    V = VIEWS[view]
    ysv = lambda v: V['base'] - v * 1000      # 縦は top/side とも「基準 − 値」で同じ形
    spanx, spany = (m1 - m0) * 1000, abs(v1 - v0) * 1000
    W = int(round(PAPER_W * DPI / 25.4 * SS))
    H = int(round(W * spany / spanx))
    k = W / spanx                             # 1車体mm あたりの px
    ox, oy = V['x0'] + m0 * 1000, ysv(v0)
    PXx = lambda x: (x - ox) * k
    PXy = lambda y: (y - oy) * k
    X = lambda m: PXx(V['x0'] + m * 1000)
    Y = lambda v: PXy(ysv(v))
    pw = lambda mm: max(1, int(round(mm * DPI / 25.4 * SS)))
    fs = lambda mm: max(8, int(round(mm * DPI / 25.4 * SS)))

    im = Image.new('RGB', (W, H), (255, 255, 255))
    dr = ImageDraw.Draw(im)
    f_num = ImageFont.truetype(FONTP, fs(2.4))
    f_zone = ImageFont.truetype(FONTP, fs(3.2))

    def label(x, y, t, font, col):
        """罫線や車体線に文字が潰されないよう、下に白地を敷いてから書く。"""
        b = dr.textbbox((x, y), t, font=font, anchor='mm')
        # 左右は広めに空ける＝罫線が文字に触れていると「-1.0」のように読めてしまう
        pad = fs(2.4) * .55
        dr.rectangle([b[0] - pad, b[1] - pad * .2, b[2] + pad, b[3] + pad * .2], fill=(255, 255, 255))
        dr.text((x, y), t, font=font, fill=col, anchor='mm')

    # 0.1m方眼（0.5mごとに濃く）
    i = int(round(m0 * 10))
    while i <= int(round(m1 * 10)):
        m, major = i / 10.0, (i % 5 == 0)
        dr.line([(X(m), 0), (X(m), H)], fill=(168, 168, 168) if major else (222, 222, 222),
                width=pw(.26 if major else .12))
        if major:
            for yy in (fs(2.4) * .9, H - fs(2.4) * .9):
                label(X(m), yy, '%.1f' % m, f_num, (130, 130, 130))
        i += 1
    lo, hi = min(v0, v1), max(v0, v1)
    j = int(round(lo * 10))
    while j <= int(round(hi * 10)):
        v, major = j / 10.0, (j % 5 == 0)
        dr.line([(0, Y(v)), (W, Y(v))], fill=(168, 168, 168) if major else (222, 222, 222),
                width=pw(.26 if major else .12))
        if major:
            t = ('0' if abs(v) < 1e-9 else '%+.1f' % v) if view == 'top' else '%.1f' % v
            for xx in (fs(2.4) * 1.5, W - fs(2.4) * 1.5):
                label(xx, Y(v), t, f_num, (130, 130, 130))
        j += 1

    # 区画の仕切りと名前
    for a, b, t in ZONES:
        if b < m0 or a > m1:
            continue
        if a > m0:
            for yy in range(0, H, pw(2.2)):
                dr.line([(X(a), yy), (X(a), yy + pw(1.3))], fill=(205, 185, 145), width=pw(.3))
        label(X((max(a, m0) + min(b, m1)) / 2), fs(3.2) * 2.4, t, f_zone, (188, 167, 118))

    # top＝中心線（左右の符号の基準）／side＝地面（高さ0）
    if view == 'top':
        for xx in range(0, W, pw(3.4)):
            dr.line([(xx, Y(0)), (xx + pw(2.0), Y(0))], fill=(206, 150, 150), width=pw(.35))
    else:
        dr.line([(0, Y(0)), (W, Y(0))], fill=(170, 155, 155), width=pw(.6))

    for cls, d in PATHS[view]:
        w, c = pw(WIDTH_MM[cls]), COL[cls]
        for sp in subpaths(d):
            pts = [(PXx(x), PXy(y)) for x, y in sp]
            if len(pts) > 1:
                dr.line(pts, fill=c, width=w, joint='curve')

    im = im.resize((W // SS, H // SS), Image.LANCZOS)
    im.save(dst, optimize=True)
    return dst, W // SS, H // SS


# 置いてもらう部品。id は wiring-net.json の parts と合わせてある
PARTS = [
    ('battery', 'バッテリー'), ('quadro', 'メーター'),
    ('ign_sw', 'キースイッチ'), ('starter_sw', '始動レバー'),
    ('starter', 'スターター'), ('dynamo', 'ダイナモ'),
    ('regulator', 'レギュレータ'), ('f1', 'ヒューズ箱'),
    ('oil_sender', '油圧センダ'), ('body', '車体アース'),
    ('free1', '（自由記入1）'), ('free2', '（自由記入2）'),
    ('free3', '（自由記入3）'), ('free4', '（自由記入4）'),
]

TXT = RGBColor(0x3a, 0x35, 0x2c)
SUB = RGBColor(0x8d, 0x85, 0x74)
ACC = RGBColor(0x2f, 0x5d, 0x50)


def tb(sl, x, y, w, h, s, size=10, bold=False, color=TXT):
    t = sl.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h)).text_frame
    t.word_wrap = True
    t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
    for i, line in enumerate(s.split('\n')):
        p = t.paragraphs[0] if i == 0 else t.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(size), bold, color, 'Meiryo'
    return t


def head(sl, title, sub):
    tb(sl, 10, 7, 200, 8, title, 15, True)
    tb(sl, 10, 16, 270, 6, sub, 9, False, SUB)


def put_view(sl, png, view, m0, m1, v0, v1, top_mm):
    """図を置く。名前に対応表を埋める＝読み取り側は動かされても正しく換算できる。"""
    path, pw_, ph_ = png
    h = PAPER_W * ph_ / pw_
    pic = sl.shapes.add_picture(path, Mm(8), Mm(top_mm), Mm(PAPER_W), Mm(h))
    pic.name = 'CARVIEW|%s|%.4f|%.4f|%.4f|%.4f' % (view, m0, m1, v0, v1)
    return top_mm + h


def palette(sl, y0, note):
    tb(sl, 10, y0, 277, 5, note, 8.5, False, SUB)
    w, hh, gx, gy = 19.0, 7.6, 1.4, 1.8
    for i, (pid, lab) in enumerate(PARTS):
        c, r = i % 7, i // 7
        s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(10 + c * (w + gx)), Mm(y0 + 5.5 + r * (hh + gy)), Mm(w), Mm(hh))
        s.name = 'PART|%s' % pid
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0xff, 0xfd, 0xf6)
        s.line.color.rgb = SUB if pid.startswith('free') else ACC
        s.line.width = Pt(.9)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r_ = p.add_run()
        r_.text = lab
        r_.font.size, r_.font.color.rgb, r_.font.name = Pt(7.5), TXT, 'Meiryo'


if not os.path.isdir(IMGD):
    os.makedirs(IMGD)
IMG = {
    'top':  render('top', -0.05, 3.05, 0.80, -0.80, IMGD + '/top.png'),
    'side': render('side', -0.05, 3.05, 1.50, -0.10, IMGD + '/side.png'),
}

prs = Presentation()
prs.slide_width, prs.slide_height = Mm(297), Mm(210)
blank = prs.slide_layouts[6]

# 1枚目：使い方
sl = prs.slides.add_slide(blank)
head(sl, 'FIAT 500F｜部品の位置と配線の這わせ方 — 記入用',
     'A4横。図はどれも 1マス = 0.1m の実寸です。印刷して手書きしても、パワポ上で図形を動かしても構いません。')
tb(sl, 10, 30, 135, 100,
   '■ パワポ上で指定する場合（おすすめ）\n'
   '1. 各シート下の「部品」の角丸を、図の上の正しい位置へドラッグする。\n'
   '2. 大きさは変えなくてよい。位置は【図形の中央】で読み取る。\n'
   '3. 要らない部品は消す。足りない物は「自由記入」の文字を書き換えて使う。\n'
   '4. 配線は［挿入］→［図形］→［フリーフォーム：図形］で、部品から部品へ引く。\n'
   '　 端点にいちばん近い部品を自動で判定するので、線に名前は要らない。\n'
   '5. 上書き保存して .pptx のまま返す。座標をそのまま数値で読み取る。\n'
   '※ 図が小さければ Ctrl＋マウスホイールで画面を拡大してください（拡大シートは要りません）。', 10)
tb(sl, 152, 30, 135, 100,
   '■ 紙に手書きする場合\n'
   '・そのまま A4横で印刷する。方眼の数字を読めば位置が取れるので、\n'
   '　印刷の拡大率がずれていても問題ない。\n\n'
   '■ 座標の決まり\n'
   '・前後＝いつも「車の前端から何m」。前端が 0.0、後端が 2.97。\n'
   '・上から見た図の縦＝「中心線から左右に何m」。\n'
   '　＋が車の右＝紙の上／−が車の左＝紙の下。\n'
   '・横から見た図の縦＝「地面から何m」。下の太い線が高さ 0。\n'
   '・横から見た図は【車の左側面】を見ている。ノーズはどちらも左。\n\n'
   '■ 区画の目安\n'
   '・前トランク 0〜1.02m／室内 1.02〜2.05m／エンジンルーム 2.05〜2.97m', 10)
tb(sl, 10, 190, 277, 8,
   '※ この図は 3Dモデルの正射投影から起こした実寸線画。部品の位置と配線のルートだけが'
   '分かればよく、絵の細部は気にしなくて構いません。', 8.5, False, SUB)


def sheet(title, sub, img, view, m0, m1, v0, v1):
    """1つの面図につき1枚。部品と配線を同じ紙の上でやる＝配線を引くとき部品が見えている。"""
    sl = prs.slides.add_slide(blank)
    head(sl, title, sub)
    y = put_view(sl, img, view, m0, m1, v0, v1, 23) + 2.5
    palette(sl, y,
            '▼ ここの角丸を図の上へドラッグしてください（位置＝図形の中央）。'
            '要らない物は削除、足りない物は「自由記入」の文字を書き換えて使ってください。')
    tb(sl, 10, y + 24.5, 277, 10,
       '▼ 配線は［挿入］→［図形］→［フリーフォーム：図形］で、部品から部品へ引いてください'
       '（折れ点をクリックしていき、最後にダブルクリック）。線に名前は要りません――'
       '端点にいちばん近い部品を自動で判定します。', 9)


sheet('① 上から見た位置', '部品の前後・左右と、配線の這わせ方（全体）', IMG['top'], 'top', -0.05, 3.05, 0.80, -0.80)
sheet('② 横から見た高さ', '部品の高さと、配線の這わせ方（全体・車の左側面）', IMG['side'], 'side', -0.05, 3.05, 1.50, -0.10)

prs.save(DST)
print('%s (%.0f KB / %d 枚)' % (DST, os.path.getsize(DST) / 1024, len(prs.slides._sldIdLst)))
